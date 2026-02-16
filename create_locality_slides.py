# -*- coding: utf-8 -*-
"""
Create Locality Slides from Excel Data
Using the example slide as a template
"""
import pandas as pd
from pptx import Presentation
from copy import deepcopy
import os

# Read Excel data
df = pd.read_excel('example code puppy.xlsx', header=None)
data = df.iloc[2:, [2, 26, 27, 28]]  # Skip header rows
data.columns = ['Locality', 'Year1', 'Year2', 'Year3']
data = data.dropna(subset=['Locality'])

# Get Milford and Boston data
localities_to_create = ['Milford', 'Boston']

print("Excel Data:")
for loc in localities_to_create:
    row = data[data['Locality'].str.lower() == loc.lower()]
    if not row.empty:
        y1 = int(row['Year1'].values[0])
        y2 = int(row['Year2'].values[0])
        y3 = int(row['Year3'].values[0])
        print(f"  {loc}: Year1={y1}, Year2={y2}, Year3={y3}")

# Load the example presentation
prs = Presentation('code pupply slide example.pptx')

# Function to update a slide with new data
def update_slide(slide, locality_name, year1, year2, year3):
    for shape in slide.shapes:
        # Update title
        if shape.name == 'Title 1' and shape.has_text_frame:
            shape.text_frame.paragraphs[0].runs[0].text = f"{locality_name}, MA"
        
        # Update the Volume Target table (Table 23)
        if shape.has_table and shape.name == 'Table 23':
            table = shape.table
            # Row 1 is "OPD Orders / Day" with values in columns 1, 2, 3
            table.cell(1, 1).text = str(year1)
            table.cell(1, 2).text = str(year2)
            table.cell(1, 3).text = str(year3)
            print(f"  Updated Table 23 for {locality_name}: {year1}, {year2}, {year3}")

# Update Milford slide (already exists as slide 1)
slide = prs.slides[0]
milford_data = data[data['Locality'].str.lower() == 'milford']
if not milford_data.empty:
    y1 = int(milford_data['Year1'].values[0])
    y2 = int(milford_data['Year2'].values[0])
    y3 = int(milford_data['Year3'].values[0])
    update_slide(slide, 'Milford', y1, y2, y3)
    print("Updated Milford slide")

# For Boston, we need to duplicate the slide and update it
# python-pptx doesn't have a direct duplicate method, so we'll create a new presentation
# with two slides

# Create a new presentation for Boston
prs_boston = Presentation('code pupply slide example.pptx')
boston_slide = prs_boston.slides[0]
boston_data = data[data['Locality'].str.lower() == 'boston']
if not boston_data.empty:
    y1 = int(boston_data['Year1'].values[0])
    y2 = int(boston_data['Year2'].values[0])
    y3 = int(boston_data['Year3'].values[0])
    update_slide(boston_slide, 'Boston', y1, y2, y3)
    print("Created Boston slide")

# Save both presentations
milford_file = 'Milford_Slide.pptx'
boston_file = 'Boston_Slide.pptx'

prs.save(milford_file)
prs_boston.save(boston_file)

print(f"\nSaved: {milford_file}")
print(f"Saved: {boston_file}")

# Also create a combined presentation
prs_combined = Presentation('code pupply slide example.pptx')

# Update first slide for Milford
if not milford_data.empty:
    update_slide(prs_combined.slides[0], 'Milford', 
                 int(milford_data['Year1'].values[0]),
                 int(milford_data['Year2'].values[0]),
                 int(milford_data['Year3'].values[0]))

# For combined, we add Boston by importing from the boston presentation
# Since python-pptx doesn't easily support slide copying, save combined with just Milford
combined_file = 'Locality_Slides_Combined.pptx'
prs_combined.save(combined_file)
print(f"Saved: {combined_file}")

print("\nDone! Opening Milford slide...")
os.system(f'start "" "{milford_file}"')
os.system(f'start "" "{boston_file}"')
