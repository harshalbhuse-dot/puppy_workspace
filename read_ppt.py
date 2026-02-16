from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation('code pupply slide example.pptx')

for i, slide in enumerate(prs.slides):
    print(f'=== Slide {i+1} ===')
    for j, shape in enumerate(slide.shapes):
        print(f'\n  [{j}] {shape.name} (type: {shape.shape_type})')
        
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(f'      Text: {text[:100]}')
        
        if shape.has_table:
            table = shape.table
            print(f'      TABLE: {len(table.rows)}x{len(table.columns)}')
            for row_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace('\n', ' ')[:30] for cell in row.cells]
                print(f'        Row {row_idx}: {cells}')
